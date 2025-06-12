# ppt_builder.py
from __future__ import annotations

import logging
import re
import sys
import tempfile
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.presentation import Presentation as _PresClass
from pptx.util import Pt

try:
    from .extractor import extract_pictures
    from .image_utils import unsplash_search
except ImportError:  # 脚本测试
    from extractor import extract_pictures          # type: ignore
    from image_utils import unsplash_search         # type: ignore

__all__ = ["prepare_image_candidates", "render_ppt"]
LOG = logging.getLogger(__name__)

# ───────────────── 进度反馈 ─────────────────
def _default_progress(pct: float, txt: str) -> None:
    bar = int(pct * 30)
    sys.stdout.write(f"\r[{u'█' * bar:<30}] {pct*100:5.1f}%  {txt:<28}")
    sys.stdout.flush()
    if pct >= 1.0:
        sys.stdout.write("\n")


def _touch_ui_event_loop() -> None:
    try:
        from PyQt5.QtWidgets import QApplication  # type: ignore
        app = QApplication.instance()
        if app:
            app.processEvents()
    except ImportError:
        pass


def _report(cb: Optional[Callable[[float, str], None]], pct: float, msg: str):
    (cb or _default_progress)(pct, msg)
    _touch_ui_event_loop()

# ───────────────── 基础工具 ─────────────────
def _ensure_presentation(src: Union[str, Path, object]) -> _PresClass:
    if isinstance(src, (str, Path)):
        return Presentation(str(src))
    if isinstance(src, _PresClass):
        return src
    if hasattr(src, "slides"):
        return src
    raise TypeError("必须是文件路径或 python-pptx Presentation 对象")


def _replace_picture(slide, r_id: str, new_blob: bytes) -> None:
    slide.part.rels[r_id].target.part.blob = new_blob


# ───────────────── Markdown ↔ Slide ─────────────────
_md_slide_pat = re.compile(r"^#{1,6}\s*slide\s*[-_ ]?\s*\d+\s*$", re.I | re.M)

def _split_markdown_slides(md: str) -> List[str]:
    """把整份 Markdown 切成每页字符串"""
    md = md.strip("\n")
    if not md:
        return []

    if re.search(r"\n?---+\n", md):
        pages = re.split(r"\n?---+\n+", md)
        return [p.strip() for p in pages if p.strip()]

    matches = list(_md_slide_pat.finditer(md))
    if not matches:
        return [md.strip()]

    pages: List[str] = []
    for i, m in enumerate(matches):
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(md)
        body = md[start:end].strip()
        if body:
            pages.append(body)
    return pages


_bullet_pat = re.compile(r"^\s*([-*]|[0-9]+[.)])\s+")
_placeholder_only_pat = re.compile(
    r"^\s*(?:[-*]|[0-9]+[.)]|•)?\s*\{\{S\d+_P\d+\}\}\s*(?:<[^>]+>\s*)?$",
    re.I,
)

def _clean_md_line(line: str) -> str:
    """去掉标题/列表标记，统一 bullet 前缀为 '• '"""
    line = line.strip()
    if line.startswith("#"):
        return line.lstrip("#").strip()
    m = _bullet_pat.match(line)
    if m:
        return "• " + line[m.end():].strip()
    return line

def _is_placeholder_only(text: str) -> bool:
    return bool(_placeholder_only_pat.match(text))


# ───────────────── 文字写入 ─────────────────
def _write_page_to_slide(slide, lines: List[str]) -> None:
    """把 Markdown 行写入 slide，遇到占位符行只占位不写字"""
    if not lines:
        return

    # 获取模板中的标题框、正文框
    title_tf = getattr(slide.shapes, "title", None)
    body_frames = [
        sh.text_frame
        for sh in slide.shapes
        if getattr(sh, "has_text_frame", False)
        and (title_tf is None or sh != title_tf)
    ]
    if not body_frames:  # 模板无正文框 → 创建一个
        body_frames.append(
            slide.shapes.add_textbox(Pt(40), Pt(100), Pt(860), Pt(400)).text_frame
        )

    # 清空正文框
    for tf in body_frames:
        tf.clear()

    # 处理标题行
    first_line = _clean_md_line(lines[0])
    if not _is_placeholder_only(first_line):
        if title_tf:
            title_tf.text = first_line
        else:  # 模板无标题框 → 动态创建
            tf = slide.shapes.add_textbox(Pt(30), Pt(20), Pt(860), Pt(60)).text_frame
            tf.text = first_line
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # 处理正文行
    frame_idx = 0
    for raw in lines[1:]:
        cleaned = _clean_md_line(raw)
        if _is_placeholder_only(cleaned):
            frame_idx += 1       # 仅占位，推进到下一文本框
            continue

        if frame_idx >= len(body_frames):   # 模板没有对应文本框 → 跳过写入
            frame_idx += 1
            continue

        tf = body_frames[frame_idx]
        para = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        para.text = cleaned
        para.font.size = Pt(18)
        frame_idx += 1


# ───────────────── prepare_image_candidates ─────────────────
def prepare_image_candidates(
    prs_or_path: Union[str, Path, object],
    keyword: str,
    *,
    per_page: int = 4,
    progress: Optional[Callable[[float, str], None]] = None,
) -> Dict[Tuple[int, str], Dict[str, object]]:
    prs = _ensure_presentation(prs_or_path)
    ret: Dict[Tuple[int, str], Dict[str, object]] = {}

    pics = list(extract_pictures(prs))
    total = len(pics) or 1
    for i, (s_idx, r_id, blob) in enumerate(pics, 1):
        _report(progress, i / total * 0.9, f"下载候选图 {i}/{total}")
        try:
            cand = unsplash_search(keyword, per_page=per_page)
        except Exception as exc:  # noqa: BLE001
            LOG.warning("获取候选图失败 (%s-%s): %s", s_idx, r_id, exc)
            cand = []
        ret[(s_idx, r_id)] = {"origin": blob, "candidates": cand}

    _report(progress, 1.0, "图片候选准备完成")
    return ret


# ───────────────── render_ppt ─────────────────
def render_ppt(
    template_path: Path | str,
    markdown: str,
    topic: str,
    user_choices: Dict[Tuple[int, str], Optional[bytes]],
    *,
    progress: Optional[Callable[[float, str], None]] = None,
) -> Path:
    prs = _ensure_presentation(template_path)
    blank_layout = prs.slide_layouts[1]

    pages = _split_markdown_slides(markdown)
    total_page = max(len(pages), 1)

    for idx, md in enumerate(pages, 1):
        # 逐行切分，列表里每个元素就是一行
        lines = [l for l in re.split(r"\n+", md) if l.strip()]
        slide = (
            prs.slides[idx - 1]
            if idx - 1 < len(prs.slides)
            else prs.slides.add_slide(blank_layout)
        )
        _write_page_to_slide(slide, lines)
        _report(progress, idx / total_page * 0.45, f"写入第 {idx}/{total_page} 页")

    # 替换用户挑选的图
    total_rep = max(len(user_choices), 1)
    for j, ((s_idx, r_id), blob) in enumerate(user_choices.items(), 1):
        if blob:
            try:
                _replace_picture(prs.slides[s_idx], r_id, blob)
            except Exception as exc:  # noqa: BLE001
                LOG.warning("替换图片失败 (%s,%s): %s", s_idx, r_id, exc)
        _report(progress, 0.45 + j / total_rep * 0.5, f"替换图片 {j}/{total_rep}")

    out = Path(tempfile.gettempdir()) / f"ai_ppt_{topic}.pptx"
    prs.save(out)
    _report(progress, 1.0, "PPT 渲染完成")
    LOG.info("PPT 渲染完成 → %s", out)
    return out
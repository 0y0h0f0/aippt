# extract_worker.py
# -*- coding: utf-8 -*-
"""
PPTX → Markdown（带占位符）＋图片提取
-----------------------------------------------------------------
外部接口
    pptx_to_markdown(path)   -> (markdown, placeholder_map)
    extract_pictures(prs)    -> [(slide_idx, rId, blob), ...]
    ExtractWorker(QThread)   -> finished / error 信号
"""
from __future__ import annotations

import traceback
from pathlib import Path
from typing import Dict, List, Tuple

from PyQt5.QtCore import QThread, pyqtSignal
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

# ───────── 常量 & 类型 ─────────
PlaceholderMap = Dict[str, Tuple[int, int]]      # "{{S1_P2}}" -> (slide_idx, shape_idx)
SHORT_LIMIT    = 10                              # ≤N 字视为“短文本” → 3 级标题

__all__ = [
    "pptx_to_markdown",
    "extract_pictures",
    "PlaceholderMap",
    "ExtractWorker",
]

# ───────── 私有工具 ─────────
def _make_ph_name(slide_idx: int, shape_idx: int) -> str:
    """生成形如 {{S1_P2}} 的占位符文本"""
    return f"{{{{S{slide_idx + 1}_P{shape_idx + 1}}}}}"


def _paragraph_align_tag(para) -> str:
    """返回段落对齐 HTML 标签，仅 center / right 用于简单排版"""
    if para.alignment == PP_ALIGN.CENTER:
        return "<div align='center'>"
    if para.alignment == PP_ALIGN.RIGHT:
        return "<div align='right'>"
    return ""


# ───────── PPTX → Markdown ─────────
def pptx_to_markdown(ppt_path: Path) -> Tuple[str, PlaceholderMap]:
    """
    读取 PPT，生成带占位符的 Markdown，并返回
    占位符 → (slide_idx, shape_idx) 的映射
    """
    prs = Presentation(str(ppt_path))
    md_lines: List[str] = []
    ph_map: PlaceholderMap = {}

    # 用于在终端输出整体框架
    debug_lines: List[str] = []

    for s_idx, slide in enumerate(prs.slides):
        md_lines.append(f"## Slide {s_idx + 1}")
        debug_lines.append(f"[EXTRACT] Slide {s_idx + 1}")

        # ---------- 收集本页文本框 ----------
        items: List[Tuple[str, int, str]] = []        # (ph, char_len, align_tag)
        for sh_idx, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue

            ph = _make_ph_name(s_idx, sh_idx)
            ph_map[ph] = (s_idx, sh_idx)

            # 直接取整个 text_frame 的文本，避免 run 遗漏
            raw_text = shape.text_frame.text.replace("\n", "")
            ln = len(raw_text.strip())

            # 取首段设置对齐标签
            align_tag = ""
            if shape.text_frame.paragraphs:
                align_tag = _paragraph_align_tag(shape.text_frame.paragraphs[0])

            items.append((ph, ln, align_tag))
            debug_lines.append(f"  · {ph} (len:{ln})")

        # ---------- 按长度输出 Markdown ----------
        if not items:
            md_lines.append("(No text on this slide)\n")
            continue

        i = 0
        while i < len(items):
            ph, ln, tag = items[i]

            if ln <= SHORT_LIMIT:                           # 短文本 → 三级标题
                block: List[str] = []
                j = i
                while j < len(items) and items[j][1] <= SHORT_LIMIT:
                    ph_j, ln_j, tag_j = items[j]
                    seg = ph_j
                    if tag_j:
                        seg += f" {tag_j}"
                    seg += f" <!--len:{ln_j}-->"
                    block.append(seg)
                    j += 1
                md_lines.extend([f"### {seg}" for seg in block])
                i = j
            else:                                           # 普通文本框
                line = f"- {ph}"
                if tag:
                    line += f" {tag}"
                line += f" <!--len:{ln}-->"
                md_lines.append(line)
                i += 1

        md_lines.append("")        # 页尾空行

    # ---------- 输出调试框架 ----------
    print("\n".join(debug_lines))

    return "\n".join(md_lines), ph_map


# ───────── 图片提取 ─────────
def extract_pictures(prs: Presentation) -> List[Tuple[int, str, bytes]]:
    """提取所有图片 (slide_idx, rId, blob)。重复引用会返回多条记录"""
    pictures: List[Tuple[int, str, bytes]] = []
    for s_idx, slide in enumerate(prs.slides):
        for rel in slide.part.rels.values():
            if rel.reltype.endswith("/image"):
                pictures.append((s_idx, rel.rId, rel.target_part.blob))
    return pictures


# ───────── 线程封装 ─────────
class ExtractWorker(QThread):
    """
    异步线程：解析 PPT → Markdown
    finished(markdown: str, error: object)  # 正常 error=None
    """
    finished = pyqtSignal(str, object)

    def __init__(self, ppt_path: Path):
        super().__init__()
        self._ppt_path = ppt_path

    def run(self) -> None:
        try:
            md, _ = pptx_to_markdown(self._ppt_path)
            self.finished.emit(md, None)
        except Exception as exc:
            traceback.print_exc()
            self.finished.emit("", exc)